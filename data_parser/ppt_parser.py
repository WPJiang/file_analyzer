import os
import hashlib
import subprocess
import tempfile
from typing import List
from .base_parser import BaseParser, DataBlock, ModalityType


class PPTParser(BaseParser):
    """PPT文件解析器

    将PPT文件解析为数据块，每页合并为一个数据块。
    """

    def __init__(self, cache_dir: str = "cache/data_blocks"):
        super().__init__()
        self.supported_extensions = ['pptx', 'ppt']
        self.cache_dir = cache_dir

    def _get_cache_path(self, file_path: str) -> str:
        """生成缓存目录路径"""
        file_hash = hashlib.md5(os.path.abspath(file_path).encode()).hexdigest()[:12]
        file_name = os.path.splitext(os.path.basename(file_path))[0]

        cache_path = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            self.cache_dir,
            f"{file_name}_{file_hash}"
        )

        os.makedirs(cache_path, exist_ok=True)
        return cache_path

    def parse(self, file_path: str) -> List[DataBlock]:
        """解析PPT文件"""
        if not self.can_parse(file_path):
            raise ValueError(f"Unsupported file format: {file_path}")

        ext = file_path.lower().split('.')[-1]

        if ext == 'ppt':
            return self._parse_ppt(file_path)
        else:
            return self._parse_pptx(file_path)

    def _parse_pptx(self, file_path: str) -> List[DataBlock]:
        """解析pptx文件 - 每页合并为一个数据块"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint file parsing. "
                "Install it with: pip install python-pptx"
            )

        blocks = []
        try:
            prs = Presentation(file_path)
        except Exception as e:
            # 文件打开失败，返回空列表，让上层处理
            print(f"[PPTParser] 无法打开PPT文件: {e}")
            return blocks

        # 获取缓存目录
        cache_path = self._get_cache_path(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            slide_tables = []

            # 收集该页所有文本和表格
            for shape in slide.shapes:
                # 解析文本
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

                # 解析表格
                if shape.has_table:
                    try:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        table_text = self._table_to_text(table_data)
                        slide_tables.append(table_text)
                    except Exception as e:
                        print(f"[PPTParser] 解析表格失败: {e}")

            # 合并该页所有内容
            all_content = []
            if slide_texts:
                all_content.append("\n\n".join(slide_texts))
            if slide_tables:
                all_content.append("\n\n".join(slide_tables))

            if all_content:
                combined_text = "\n\n---\n\n".join(all_content)
                text_cache_path = os.path.join(cache_path, f"slide_{slide_num}.txt")
                try:
                    with open(text_cache_path, 'w', encoding='utf-8') as f:
                        f.write(combined_text)
                except Exception as e:
                    print(f"[PPTParser] 保存文本缓存失败: {e}")
                    text_cache_path = None

                block = DataBlock(
                    block_id=self._generate_block_id(file_path, slide_num - 1),
                    modality=ModalityType.TEXT,
                    addr=text_cache_path,
                    file_path=file_path,
                    page_number=slide_num,
                    metadata={
                        'source': 'python-pptx',
                        'slide_number': slide_num,
                        'text_length': len(combined_text)
                    }
                )
                blocks.append(block)

        return blocks

    def _parse_ppt(self, file_path: str) -> List[DataBlock]:
        """解析ppt文件（需要LibreOffice转换）"""
        temp_dir = tempfile.gettempdir()
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        pptx_path = os.path.join(temp_dir, f"{base_name}.pptx")

        try:
            # 尝试使用LibreOffice转换
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pptx',
                 '--outdir', temp_dir, file_path],
                capture_output=True,
                text=True,
                timeout=60  # 添加超时限制
            )

            if result.returncode == 0 and os.path.exists(pptx_path):
                blocks = self._parse_pptx(pptx_path)
                try:
                    os.remove(pptx_path)
                except Exception:
                    pass
                return blocks
            else:
                error_msg = result.stderr if result.stderr else "未知错误"
                print(f"[PPTParser] LibreOffice转换失败: {error_msg}")
                # 返回空列表而不是抛出异常，让文件仍能被处理
                return []

        except subprocess.TimeoutExpired:
            print("[PPTParser] LibreOffice转换超时")
            return []
        except FileNotFoundError:
            print("[PPTParser] 未找到LibreOffice (soffice命令)，请安装LibreOffice或将.ppt转换为.pptx格式")
            return []
        except Exception as e:
            print(f"[PPTParser] 解析.ppt文件失败: {e}")
            return []

    def _table_to_text(self, table_data: List[List[str]]) -> str:
        """将表格转换为文本"""
        rows = []
        for row in table_data:
            row_text = ' | '.join(str(cell) if cell else '' for cell in row)
            rows.append(row_text)
        return '\n'.join(rows)