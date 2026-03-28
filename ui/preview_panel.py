import os
import io
import tempfile
import zipfile
from typing import Optional, List
from PyQt5.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel, QTextEdit,
    QStackedWidget, QScrollArea, QFrame, QMessageBox,
    QPushButton, QProgressBar, QSpinBox, QSlider, QComboBox,
    QTableWidget, QTableWidgetItem, QTextBrowser, QSplitter
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QSize
from PyQt5.QtGui import QPixmap, QImage, QFont, QTextCursor

from .utils import get_font_sizes, get_window_sizes, get_icon_sizes

# 导入 PIL 用于更多图片格式支持
try:
    from PIL import Image
    PIL_SUPPORT = True
except ImportError:
    PIL_SUPPORT = False

# 尝试导入 pillow-heif 用于 HEIC/HEIF 格式支持
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
    HEIF_SUPPORT = True
except ImportError:
    HEIF_SUPPORT = False

# 导入文档处理库
try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False


class PDFPreviewWidget(QWidget):
    """PDF预览组件 - 支持翻页和缩放"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.doc = None
        self.current_page = 0
        self.total_pages = 0
        self.zoom_level = 1.0
        self.scale_factor = 1.5  # 基础缩放比例
        self.font_sizes = get_font_sizes()
        self.icon_sizes = get_icon_sizes()
        self.window_sizes = get_window_sizes()

        self.init_ui()

    def init_ui(self):
        """初始化UI"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(self.window_sizes['spacing_small'])

        # 工具栏
        toolbar = QWidget()
        toolbar_layout = QHBoxLayout(toolbar)
        toolbar_layout.setContentsMargins(
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small']
        )
        toolbar_layout.setSpacing(self.window_sizes['spacing_normal'])

        # 上一页按钮
        self.prev_btn = QPushButton("◀ 上一页")
        self.prev_btn.setMinimumHeight(self.window_sizes['button_height'])
        self.prev_btn.setStyleSheet(self.get_button_style("#2196F3"))
        self.prev_btn.clicked.connect(self.prev_page)
        toolbar_layout.addWidget(self.prev_btn)

        # 页码显示
        self.page_label = QLabel("第 0 / 0 页")
        self.page_label.setStyleSheet(f"font-size: {self.font_sizes['normal']}px; color: #333;")
        toolbar_layout.addWidget(self.page_label)

        # 下一页按钮
        self.next_btn = QPushButton("下一页 ▶")
        self.next_btn.setMinimumHeight(self.window_sizes['button_height'])
        self.next_btn.setStyleSheet(self.get_button_style("#2196F3"))
        self.next_btn.clicked.connect(self.next_page)
        toolbar_layout.addWidget(self.next_btn)

        toolbar_layout.addStretch()

        # 缩放控制
        zoom_label = QLabel("缩放:")
        zoom_label.setStyleSheet("color: #666;")
        toolbar_layout.addWidget(zoom_label)

        self.zoom_combo = QComboBox()
        self.zoom_combo.addItems(["50%", "75%", "100%", "125%", "150%", "200%", "300%"])
        self.zoom_combo.setCurrentText("100%")
        self.zoom_combo.currentTextChanged.connect(self.on_zoom_changed)
        self.zoom_combo.setMinimumHeight(self.window_sizes['input_height'])
        self.zoom_combo.setStyleSheet(f"""
            QComboBox {{
                padding: 4px 8px;
                border: 1px solid #ddd;
                border-radius: 4px;
                min-width: 80px;
                font-size: {self.font_sizes['small']}px;
            }}
        """)
        toolbar_layout.addWidget(self.zoom_combo)

        # 缩小按钮
        zoom_out_btn = QPushButton("-")
        btn_size = self.icon_sizes['small']
        zoom_out_btn.setFixedSize(btn_size, btn_size)
        zoom_out_btn.setStyleSheet(self.get_button_style("#607D8B"))
        zoom_out_btn.clicked.connect(self.zoom_out)
        toolbar_layout.addWidget(zoom_out_btn)

        # 放大按钮
        zoom_in_btn = QPushButton("+")
        zoom_in_btn.setFixedSize(btn_size, btn_size)
        zoom_in_btn.setStyleSheet(self.get_button_style("#607D8B"))
        zoom_in_btn.clicked.connect(self.zoom_in)
        toolbar_layout.addWidget(zoom_in_btn)

        layout.addWidget(toolbar)

        # 预览区域
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setAlignment(Qt.AlignCenter)
        scroll.setStyleSheet("""
            QScrollArea {
                border: 1px solid #ddd;
                background-color: #f5f5f5;
            }
        """)

        self.image_label = QLabel()
        self.image_label.setAlignment(Qt.AlignCenter)
        self.image_label.setStyleSheet("background-color: white;")

        scroll.setWidget(self.image_label)
        layout.addWidget(scroll)

        # 页面信息
        self.info_label = QLabel("")
        self.info_label.setStyleSheet("color: #666; padding: 5px;")
        self.info_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(self.info_label)

    def get_button_style(self, color: str) -> str:
        """获取按钮样式"""
        return f"""
            QPushButton {{
                background-color: {color};
                color: white;
                border: none;
                padding: 6px 12px;
                border-radius: 4px;
                font-size: {self.font_sizes['button']}px;
            }}
            QPushButton:hover {{
                background-color: {color}dd;
            }}
            QPushButton:disabled {{
                background-color: #ccc;
            }}
        """
    
    def load_pdf(self, file_path: str) -> bool:
        """加载PDF文件"""
        if not PDF_SUPPORT:
            return False
        
        try:
            self.doc = fitz.open(file_path)
            self.total_pages = len(self.doc)
            self.current_page = 0
            self.zoom_level = 1.0
            self.zoom_combo.setCurrentText("100%")
            self.update_page()
            return True
        except Exception as e:
            print(f"加载PDF失败: {e}")
            return False
    
    def update_page(self):
        """更新当前页面显示"""
        if not self.doc or self.total_pages == 0:
            return
        
        # 更新页码标签
        self.page_label.setText(f"第 {self.current_page + 1} / {self.total_pages} 页")
        
        # 更新按钮状态
        self.prev_btn.setEnabled(self.current_page > 0)
        self.next_btn.setEnabled(self.current_page < self.total_pages - 1)
        
        # 渲染页面
        try:
            page = self.doc[self.current_page]
            
            # 计算缩放矩阵
            mat = fitz.Matrix(self.scale_factor * self.zoom_level, 
                             self.scale_factor * self.zoom_level)
            pix = page.get_pixmap(matrix=mat)
            
            # 转换为QImage
            img = QImage(pix.samples, pix.width, pix.height, 
                        pix.stride, QImage.Format_RGB888)
            pixmap = QPixmap.fromImage(img)
            
            self.image_label.setPixmap(pixmap)
            self.image_label.setFixedSize(pixmap.size())
            
            # 更新信息
            self.info_label.setText(
                f"页面尺寸: {page.rect.width:.0f} x {page.rect.height:.0f} pt | "
                f"缩放: {self.zoom_level * 100:.0f}%"
            )
        except Exception as e:
            self.info_label.setText(f"渲染页面失败: {e}")
    
    def prev_page(self):
        """上一页"""
        if self.current_page > 0:
            self.current_page -= 1
            self.update_page()
    
    def next_page(self):
        """下一页"""
        if self.current_page < self.total_pages - 1:
            self.current_page += 1
            self.update_page()
    
    def zoom_in(self):
        """放大"""
        if self.zoom_level < 3.0:
            self.zoom_level += 0.25
            self.update_zoom_combo()
            self.update_page()
    
    def zoom_out(self):
        """缩小"""
        if self.zoom_level > 0.5:
            self.zoom_level -= 0.25
            self.update_zoom_combo()
            self.update_page()
    
    def on_zoom_changed(self, text: str):
        """缩放比例改变"""
        try:
            zoom = int(text.replace('%', ''))
            self.zoom_level = zoom / 100.0
            self.update_page()
        except:
            pass
    
    def update_zoom_combo(self):
        """更新缩放下拉框"""
        zoom_text = f"{int(self.zoom_level * 100)}%"
        index = self.zoom_combo.findText(zoom_text)
        if index >= 0:
            self.zoom_combo.setCurrentIndex(index)
    
    def close_document(self):
        """关闭文档"""
        if self.doc:
            self.doc.close()
            self.doc = None


class WordPreviewWidget(QWidget):
    """Word文档预览组件"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.font_sizes = get_font_sizes()
        self.icon_sizes = get_icon_sizes()
        self.window_sizes = get_window_sizes()
        self.init_ui()

    def init_ui(self):
        """初始化UI"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(self.window_sizes['spacing_small'])

        # 工具栏
        toolbar = QWidget()
        toolbar_layout = QHBoxLayout(toolbar)
        toolbar_layout.setContentsMargins(
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small']
        )

        info_label = QLabel("Word 文档预览")
        info_label.setStyleSheet("font-weight: bold; color: #333;")
        toolbar_layout.addWidget(info_label)
        toolbar_layout.addStretch()

        layout.addWidget(toolbar)

        # 文本显示区域
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setStyleSheet("""
            QScrollArea {
                border: 1px solid #ddd;
                background-color: white;
            }
        """)

        self.text_browser = QTextBrowser()
        self.text_browser.setOpenExternalLinks(True)
        self.text_browser.setStyleSheet(f"""
            QTextBrowser {{
                border: none;
                padding: 20px;
                background-color: white;
                font-family: 'Microsoft YaHei', 'SimSun', serif;
                font-size: {self.font_sizes['normal']}pt;
                line-height: 1.6;
            }}
        """)
        
        scroll.setWidget(self.text_browser)
        layout.addWidget(scroll)
    
    def load_document(self, file_path: str) -> bool:
        """加载Word文档"""
        if not DOCX_SUPPORT:
            return False
        
        try:
            doc = Document(file_path)
            html_content = self.convert_docx_to_html(doc)
            self.text_browser.setHtml(html_content)
            return True
        except Exception as e:
            self.text_browser.setPlainText(f"无法加载Word文档: {e}")
            return False
    
    def convert_docx_to_html(self, doc) -> str:
        """将docx转换为HTML以保留样式"""
        html_parts = []
        html_parts.append("""
        <html>
        <head>
        <style>
            body { font-family: 'Microsoft YaHei', sans-serif; line-height: 1.6; }
            p { margin: 10px 0; }
            h1, h2, h3, h4, h5, h6 { margin: 15px 0 10px 0; color: #333; }
            table { border-collapse: collapse; margin: 10px 0; }
            td, th { border: 1px solid #ddd; padding: 8px; }
            ul, ol { margin: 10px 0; padding-left: 30px; }
        </style>
        </head>
        <body>
        """)
        
        for para in doc.paragraphs:
            if para.text.strip():
                # 获取段落样式
                style_name = para.style.name if para.style else 'Normal'
                
                if 'Heading' in style_name:
                    level = style_name.replace('Heading ', '')
                    try:
                        level = int(level)
                        html_parts.append(f"<h{level}>{self.escape_html(para.text)}</h{level}>")
                    except:
                        html_parts.append(f"<p>{self.escape_html(para.text)}</p>")
                else:
                    # 处理文本格式
                    text_parts = []
                    for run in para.runs:
                        text = self.escape_html(run.text)
                        if run.bold:
                            text = f"<b>{text}</b>"
                        if run.italic:
                            text = f"<i>{text}</i>"
                        if run.underline:
                            text = f"<u>{text}</u>"
                        text_parts.append(text)
                    
                    html_parts.append(f"<p>{''.join(text_parts)}</p>")
        
        # 处理表格
        for table in doc.tables:
            html_parts.append("<table>")
            for row in table.rows:
                html_parts.append("<tr>")
                for cell in row.cells:
                    html_parts.append(f"<td>{self.escape_html(cell.text)}</td>")
                html_parts.append("</tr>")
            html_parts.append("</table>")
        
        html_parts.append("</body></html>")
        return ''.join(html_parts)
    
    def escape_html(self, text: str) -> str:
        """转义HTML特殊字符"""
        return (text
                .replace('&', '&amp;')
                .replace('<', '&lt;')
                .replace('>', '&gt;')
                .replace('"', '&quot;'))


class ExcelPreviewWidget(QWidget):
    """Excel预览组件"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.workbook = None
        self.current_sheet = 0
        self.font_sizes = get_font_sizes()
        self.icon_sizes = get_icon_sizes()
        self.window_sizes = get_window_sizes()
        self.init_ui()

    def init_ui(self):
        """初始化UI"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(self.window_sizes['spacing_small'])

        # 工具栏
        toolbar = QWidget()
        toolbar_layout = QHBoxLayout(toolbar)
        toolbar_layout.setContentsMargins(
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small']
        )

        sheet_label = QLabel("工作表:")
        sheet_label.setStyleSheet("color: #666;")
        toolbar_layout.addWidget(sheet_label)

        self.sheet_combo = QComboBox()
        self.sheet_combo.currentIndexChanged.connect(self.on_sheet_changed)
        self.sheet_combo.setMinimumHeight(self.window_sizes['input_height'])
        self.sheet_combo.setStyleSheet(f"""
            QComboBox {{
                padding: 4px 8px;
                border: 1px solid #ddd;
                border-radius: 4px;
                min-width: 150px;
                font-size: {self.font_sizes['small']}px;
            }}
        """)
        toolbar_layout.addWidget(self.sheet_combo)

        toolbar_layout.addStretch()

        self.info_label = QLabel("")
        self.info_label.setStyleSheet("color: #666;")
        toolbar_layout.addWidget(self.info_label)

        layout.addWidget(toolbar)

        # 表格显示区域
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setStyleSheet("""
            QScrollArea {
                border: 1px solid #ddd;
                background-color: white;
            }
        """)

        self.table = QTableWidget()
        self.table.setStyleSheet(f"""
            QTableWidget {{
                border: none;
                gridline-color: #ddd;
                font-size: {self.font_sizes['small']}px;
            }}
            QHeaderView::section {{
                background-color: #f5f5f5;
                padding: 5px;
                border: 1px solid #ddd;
                font-weight: bold;
            }}
        """)
        
        scroll.setWidget(self.table)
        layout.addWidget(scroll)
    
    def load_workbook(self, file_path: str) -> bool:
        """加载Excel文件"""
        if not XLSX_SUPPORT:
            return False
        
        try:
            self.workbook = load_workbook(file_path, data_only=True, read_only=True)
            
            # 填充工作表下拉框
            self.sheet_combo.clear()
            for sheet_name in self.workbook.sheetnames:
                self.sheet_combo.addItem(sheet_name)
            
            if self.workbook.sheetnames:
                self.load_sheet(0)
            
            return True
        except Exception as e:
            print(f"加载Excel失败: {e}")
            return False
    
    def on_sheet_changed(self, index: int):
        """工作表切换"""
        if index >= 0:
            self.load_sheet(index)
    
    def load_sheet(self, index: int):
        """加载指定工作表"""
        if not self.workbook:
            return
        
        try:
            sheet_name = self.workbook.sheetnames[index]
            sheet = self.workbook[sheet_name]
            
            # 获取数据范围
            max_row = min(sheet.max_row, 1000)  # 限制行数
            max_col = min(sheet.max_column, 50)  # 限制列数
            
            self.table.clear()
            self.table.setRowCount(max_row)
            self.table.setColumnCount(max_col)
            
            # 设置表头
            headers = []
            for col in range(1, max_col + 1):
                headers.append(get_column_letter(col))
            self.table.setHorizontalHeaderLabels(headers)
            
            # 填充数据
            for row_idx, row in enumerate(sheet.iter_rows(min_row=1, max_row=max_row, 
                                                           max_col=max_col), 1):
                for col_idx, cell in enumerate(row, 0):
                    value = cell.value if cell.value is not None else ""
                    item = QTableWidgetItem(str(value))
                    self.table.setItem(row_idx - 1, col_idx, item)
            
            # 调整列宽
            self.table.resizeColumnsToContents()
            
            # 更新信息
            self.info_label.setText(f"共 {sheet.max_row} 行, {sheet.max_column} 列")
            
        except Exception as e:
            self.info_label.setText(f"加载工作表失败: {e}")


class PowerPointPreviewWidget(QWidget):
    """PowerPoint预览组件"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.prs = None
        self.current_slide = 0
        self.total_slides = 0
        self.font_sizes = get_font_sizes()
        self.icon_sizes = get_icon_sizes()
        self.window_sizes = get_window_sizes()
        self.init_ui()

    def init_ui(self):
        """初始化UI"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(self.window_sizes['spacing_small'])

        # 工具栏
        toolbar = QWidget()
        toolbar_layout = QHBoxLayout(toolbar)
        toolbar_layout.setContentsMargins(
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small'],
            self.window_sizes['margin_small']
        )
        toolbar_layout.setSpacing(self.window_sizes['spacing_normal'])

        self.prev_btn = QPushButton("◀ 上一页")
        self.prev_btn.setMinimumHeight(self.window_sizes['button_height'])
        self.prev_btn.setStyleSheet(self.get_button_style("#2196F3"))
        self.prev_btn.clicked.connect(self.prev_slide)
        toolbar_layout.addWidget(self.prev_btn)

        self.slide_label = QLabel("第 0 / 0 页")
        self.slide_label.setStyleSheet(f"font-size: {self.font_sizes['normal']}px; color: #333;")
        toolbar_layout.addWidget(self.slide_label)

        self.next_btn = QPushButton("下一页 ▶")
        self.next_btn.setMinimumHeight(self.window_sizes['button_height'])
        self.next_btn.setStyleSheet(self.get_button_style("#2196F3"))
        self.next_btn.clicked.connect(self.next_slide)
        toolbar_layout.addWidget(self.next_btn)

        toolbar_layout.addStretch()

        layout.addWidget(toolbar)

        # 幻灯片显示区域
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setAlignment(Qt.AlignCenter)
        scroll.setStyleSheet("""
            QScrollArea {
                border: 1px solid #ddd;
                background-color: #f5f5f5;
            }
        """)

        self.slide_widget = QTextBrowser()
        self.slide_widget.setOpenExternalLinks(True)
        self.slide_widget.setStyleSheet(f"""
            QTextBrowser {{
                border: none;
                padding: 20px;
                background-color: white;
                font-family: 'Microsoft YaHei', sans-serif;
                font-size: {self.font_sizes['normal']}px;
            }}
        """)

        scroll.setWidget(self.slide_widget)
        layout.addWidget(scroll)

    def get_button_style(self, color: str) -> str:
        """获取按钮样式"""
        return f"""
            QPushButton {{
                background-color: {color};
                color: white;
                border: none;
                padding: 6px 12px;
                border-radius: 4px;
                font-size: {self.font_sizes['button']}px;
            }}
            QPushButton:hover {{
                background-color: {color}dd;
            }}
            QPushButton:disabled {{
                background-color: #ccc;
            }}
        """
    
    def load_presentation(self, file_path: str) -> bool:
        """加载PPT文件"""
        if not PPTX_SUPPORT:
            return False
        
        try:
            self.prs = Presentation(file_path)
            self.total_slides = len(self.prs.slides)
            self.current_slide = 0
            self.update_slide()
            return True
        except Exception as e:
            print(f"加载PPT失败: {e}")
            return False
    
    def update_slide(self):
        """更新当前幻灯片显示"""
        if not self.prs or self.total_slides == 0:
            return
        
        self.slide_label.setText(f"第 {self.current_slide + 1} / {self.total_slides} 页")
        self.prev_btn.setEnabled(self.current_slide > 0)
        self.next_btn.setEnabled(self.current_slide < self.total_slides - 1)
        
        try:
            slide = self.prs.slides[self.current_slide]
            html_content = self.convert_slide_to_html(slide)
            self.slide_widget.setHtml(html_content)
        except Exception as e:
            self.slide_widget.setPlainText(f"渲染幻灯片失败: {e}")
    
    def convert_slide_to_html(self, slide) -> str:
        """将幻灯片转换为HTML"""
        html_parts = []
        html_parts.append("""
        <html>
        <head>
        <style>
            body { 
                font-family: 'Microsoft YaHei', sans-serif; 
                line-height: 1.6;
                padding: 20px;
            }
            .slide-title { 
                font-size: 24px; 
                font-weight: bold; 
                color: #333;
                margin-bottom: 20px;
                border-bottom: 2px solid #2196F3;
                padding-bottom: 10px;
            }
            .slide-content { font-size: 14px; }
            ul, ol { margin: 10px 0; padding-left: 30px; }
            li { margin: 5px 0; }
        </style>
        </head>
        <body>
        """)
        
        # 提取标题
        title_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and not title_text:
                title_text = shape.text
                break
        
        if title_text:
            html_parts.append(f'<div class="slide-title">{self.escape_html(title_text)}</div>')
        
        html_parts.append('<div class="slide-content">')
        
        # 提取内容
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text != title_text:
                text = self.escape_html(shape.text)
                # 简单处理列表
                if text.startswith('•') or text.startswith('-'):
                    html_parts.append(f"<ul><li>{text[1:].strip()}</li></ul>")
                else:
                    html_parts.append(f"<p>{text}</p>")
        
        html_parts.append('</div></body></html>')
        return ''.join(html_parts)
    
    def escape_html(self, text: str) -> str:
        """转义HTML特殊字符"""
        return (text
                .replace('&', '&amp;')
                .replace('<', '&lt;')
                .replace('>', '&gt;')
                .replace('"', '&quot;'))
    
    def prev_slide(self):
        """上一页"""
        if self.current_slide > 0:
            self.current_slide -= 1
            self.update_slide()
    
    def next_slide(self):
        """下一页"""
        if self.current_slide < self.total_slides - 1:
            self.current_slide += 1
            self.update_slide()


class PreviewWorker(QThread):
    """预览加载工作线程"""
    preview_ready = pyqtSignal(object, str)
    preview_error = pyqtSignal(str)

    def __init__(self, file_path: str):
        super().__init__()
        self.file_path = file_path

    def run(self):
        try:
            ext = os.path.splitext(self.file_path)[1].lower()

            # 图片文件
            if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif', '.ico']:
                self._load_standard_image()
            elif ext == '.svg':
                # SVG 尝试用 QPixmap 加载，可能失败
                pixmap = QPixmap(self.file_path)
                if not pixmap.isNull():
                    self.preview_ready.emit(pixmap, 'image')
                else:
                    self.preview_error.emit("SVG 格式需要额外支持，请使用外部程序打开")
            elif ext in ['.heic', '.heif']:
                self._load_heic_image()
            elif ext == '.livp':
                self._load_livp_image()
            # 文本文件
            elif ext in ['.txt', '.md', '.py', '.json', '.xml', '.html', '.css', '.js']:
                try:
                    with open(self.file_path, 'r', encoding='utf-8') as f:
                        content = f.read(10000)  # 限制预览大小
                    self.preview_ready.emit(content, 'text')
                except Exception as e:
                    self.preview_error.emit(f"无法读取文本: {str(e)}")
            # PDF文件
            elif ext == '.pdf':
                self.preview_ready.emit(self.file_path, 'pdf')
            # Word文档
            elif ext in ['.doc', '.docx']:
                self.preview_ready.emit(self.file_path, 'word')
            # Excel文档
            elif ext in ['.xls', '.xlsx']:
                self.preview_ready.emit(self.file_path, 'excel')
            # PowerPoint文档
            elif ext in ['.ppt', '.pptx']:
                self.preview_ready.emit(self.file_path, 'ppt')
            # 音频文件
            elif ext in ['.mp3', '.wav', '.m4a', '.flac']:
                self.preview_ready.emit(self.file_path, 'audio')
            # 快捷方式文件
            elif ext == '.lnk':
                self.preview_ready.emit(self.file_path, 'shortcut')
            # 其他文件
            else:
                self.preview_ready.emit(self.file_path, 'unsupported')

        except Exception as e:
            self.preview_error.emit(str(e))

    def _load_standard_image(self):
        """加载标准图片格式"""
        pixmap = QPixmap(self.file_path)
        if not pixmap.isNull():
            self.preview_ready.emit(pixmap, 'image')
        else:
            # 尝试用 PIL 加载
            if PIL_SUPPORT:
                try:
                    img = Image.open(self.file_path)
                    # 转换为 RGB 模式（如果是 RGBA 则保留）
                    if img.mode not in ('RGB', 'RGBA'):
                        img = img.convert('RGB')
                    # 转换为 QPixmap
                    data = img.tobytes('raw', img.mode)
                    qimage = QImage(data, img.width, img.height,
                                   QImage.Format_RGB888 if img.mode == 'RGB' else QImage.Format_RGBA8888)
                    pixmap = QPixmap.fromImage(qimage)
                    if not pixmap.isNull():
                        self.preview_ready.emit(pixmap, 'image')
                        return
                except Exception as e:
                    pass
            self.preview_error.emit("无法加载图片")

    def _load_heic_image(self):
        """加载 HEIC/HEIF 格式图片"""
        # 首先尝试用 QPixmap 加载
        pixmap = QPixmap(self.file_path)
        if not pixmap.isNull():
            self.preview_ready.emit(pixmap, 'image')
            return

        # 使用 PIL 加载（需要 pillow-heif）
        if PIL_SUPPORT:
            try:
                img = Image.open(self.file_path)
                if img.mode not in ('RGB', 'RGBA'):
                    img = img.convert('RGB')
                data = img.tobytes('raw', img.mode)
                qimage = QImage(data, img.width, img.height,
                               QImage.Format_RGB888 if img.mode == 'RGB' else QImage.Format_RGBA8888)
                pixmap = QPixmap.fromImage(qimage)
                if not pixmap.isNull():
                    self.preview_ready.emit(pixmap, 'image')
                    return
            except Exception as e:
                if not HEIF_SUPPORT:
                    self.preview_error.emit("HEIC 格式需要安装 pillow-heif 库\n请运行: pip install pillow-heif")
                    return
                self.preview_error.emit(f"无法加载 HEIC 图片: {str(e)}")
                return

        if not HEIF_SUPPORT:
            self.preview_error.emit("HEIC 格式需要安装 pillow-heif 库\n请运行: pip install pillow-heif")
        else:
            self.preview_error.emit("无法加载 HEIC 图片")

    def _load_livp_image(self):
        """加载 Live Photo (.livp) 格式

        Live Photo 是 Apple 的格式，包含一张静态图片和一个短视频。
        文件实际上是一个 ZIP 压缩包，我们需要提取其中的静态图片进行预览。
        """
        try:
            # 检查是否是 ZIP 文件
            if not zipfile.is_zipfile(self.file_path):
                self.preview_error.emit("无效的 Live Photo 文件格式")
                return

            with zipfile.ZipFile(self.file_path, 'r') as zf:
                # 列出压缩包中的文件
                file_list = zf.namelist()

                # 查找图片文件（通常是 HEIC 或 JPEG）
                image_extensions = ['.heic', '.heif', '.jpg', '.jpeg', '.png']
                image_file = None

                for f in file_list:
                    ext = os.path.splitext(f)[1].lower()
                    if ext in image_extensions:
                        image_file = f
                        break

                if not image_file:
                    self.preview_error.emit("Live Photo 中未找到图片文件")
                    return

                # 提取图片到临时文件
                with tempfile.NamedTemporaryFile(suffix=os.path.splitext(image_file)[1], delete=False) as tmp:
                    tmp.write(zf.read(image_file))
                    tmp_path = tmp.name

                try:
                    # 尝试加载提取的图片
                    if image_file.lower().endswith(('.heic', '.heif')):
                        # HEIC 格式
                        if PIL_SUPPORT:
                            try:
                                img = Image.open(tmp_path)
                                if img.mode not in ('RGB', 'RGBA'):
                                    img = img.convert('RGB')
                                data = img.tobytes('raw', img.mode)
                                qimage = QImage(data, img.width, img.height,
                                               QImage.Format_RGB888 if img.mode == 'RGB' else QImage.Format_RGBA8888)
                                pixmap = QPixmap.fromImage(qimage)
                                if not pixmap.isNull():
                                    self.preview_ready.emit(pixmap, 'image')
                                    return
                            except Exception as e:
                                print(f"[PreviewWorker] HEIC 加载失败: {e}")

                        if not HEIF_SUPPORT:
                            self.preview_error.emit("Live Photo 中的 HEIC 图片需要安装 pillow-heif 库")
                            return
                    else:
                        # 标准图片格式
                        pixmap = QPixmap(tmp_path)
                        if not pixmap.isNull():
                            self.preview_ready.emit(pixmap, 'image')
                            return

                    self.preview_error.emit("无法加载 Live Photo 中的图片")
                finally:
                    # 清理临时文件
                    try:
                        os.unlink(tmp_path)
                    except Exception:
                        pass

        except zipfile.BadZipFile:
            self.preview_error.emit("无效的 Live Photo 文件（不是有效的 ZIP 格式）")
        except Exception as e:
            self.preview_error.emit(f"加载 Live Photo 失败: {str(e)}")


class PreviewPanel(QWidget):
    """预览窗口组件"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.current_file = None
        self.current_metadata = None
        self.current_file_id = None
        self.current_category_system = None
        self.preview_worker = None
        self.font_sizes = get_font_sizes()
        self.icon_sizes = get_icon_sizes()
        self.window_sizes = get_window_sizes()

        self.init_ui()

    def init_ui(self):
        """初始化UI"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(
            self.window_sizes['margin_normal'],
            self.window_sizes['margin_normal'],
            self.window_sizes['margin_normal'],
            self.window_sizes['margin_normal']
        )
        layout.setSpacing(self.window_sizes['spacing_normal'])

        # 标题
        title_layout = QHBoxLayout()

        self.title_label = QLabel("文件预览")
        self.title_label.setStyleSheet(f"""
            font-size: {self.font_sizes['title']}px;
            font-weight: bold;
            color: #333;
            padding: 5px;
        """)
        title_layout.addWidget(self.title_label)

        title_layout.addStretch()

        # 文件信息
        self.file_info_label = QLabel("")
        self.file_info_label.setStyleSheet("color: #666;")
        title_layout.addWidget(self.file_info_label)
        
        layout.addLayout(title_layout)
        
        # 预览区域（使用堆叠窗口）
        self.stack = QStackedWidget()
        
        # 1. 空白页面
        self.empty_page = self.create_empty_page()
        self.stack.addWidget(self.empty_page)
        
        # 2. 图片预览页面
        self.image_page = self.create_image_page()
        self.stack.addWidget(self.image_page)
        
        # 3. 文本预览页面
        self.text_page = self.create_text_page()
        self.stack.addWidget(self.text_page)
        
        # 4. PDF预览页面
        self.pdf_widget = PDFPreviewWidget()
        self.stack.addWidget(self.pdf_widget)
        
        # 5. Word预览页面
        self.word_widget = WordPreviewWidget()
        self.stack.addWidget(self.word_widget)
        
        # 6. Excel预览页面
        self.excel_widget = ExcelPreviewWidget()
        self.stack.addWidget(self.excel_widget)
        
        # 7. PowerPoint预览页面
        self.ppt_widget = PowerPointPreviewWidget()
        self.stack.addWidget(self.ppt_widget)
        
        # 8. 音频预览页面
        self.audio_page = self.create_audio_page()
        self.stack.addWidget(self.audio_page)
        
        # 9. 不支持页面
        self.unsupported_page = self.create_unsupported_page()
        self.stack.addWidget(self.unsupported_page)
        
        # 10. 加载中页面
        self.loading_page = self.create_loading_page()
        self.stack.addWidget(self.loading_page)
        
        layout.addWidget(self.stack)

        # 语义块信息显示区域
        self.semantic_blocks_label = QLabel("")
        self.semantic_blocks_label.setStyleSheet(f"""
            color: #555;
            font-size: {self.font_sizes['normal']}px;
            padding: 5px;
            background-color: #f9f9f9;
            border: 1px solid #eee;
            border-radius: 4px;
        """)
        self.semantic_blocks_label.setWordWrap(True)
        self.semantic_blocks_label.setAlignment(Qt.AlignTop | Qt.AlignLeft)
        self.semantic_blocks_label.setVisible(False)  # 默认隐藏
        layout.addWidget(self.semantic_blocks_label)

        # 设置样式
        self.setStyleSheet("""
            PreviewPanel {
                background-color: white;
                border: 1px solid #ddd;
                border-radius: 8px;
            }
        """)
    
    def create_empty_page(self) -> QWidget:
        """创建空白页面"""
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setAlignment(Qt.AlignCenter)

        label = QLabel("请选择文件以预览")
        label.setStyleSheet(f"""
            color: #999;
            font-size: {self.font_sizes['icon_medium']}px;
            padding: 20px;
        """)
        label.setAlignment(Qt.AlignCenter)
        layout.addWidget(label)

        return page

    def create_image_page(self) -> QWidget:
        """创建图片预览页面"""
        page = QWidget()
        layout = QVBoxLayout(page)

        # 滚动区域
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setAlignment(Qt.AlignCenter)

        # 图片标签
        self.image_label = QLabel()
        self.image_label.setAlignment(Qt.AlignCenter)
        self.image_label.setStyleSheet("background-color: #f5f5f5;")

        scroll.setWidget(self.image_label)
        layout.addWidget(scroll)

        # 图片信息
        self.image_info_label = QLabel("")
        self.image_info_label.setAlignment(Qt.AlignCenter)
        self.image_info_label.setStyleSheet("color: #666; padding: 5px;")
        layout.addWidget(self.image_info_label)

        # Caption标签
        self.image_caption_label = QLabel("")
        self.image_caption_label.setAlignment(Qt.AlignCenter)
        self.image_caption_label.setStyleSheet(f"""
            color: #333;
            padding: 5px;
            font-size: {self.font_sizes['normal']}px;
            background-color: #f0f8ff;
            border-radius: 4px;
            margin: 2px;
        """)
        self.image_caption_label.setWordWrap(True)
        self.image_caption_label.hide()
        layout.addWidget(self.image_caption_label)

        # Tags标签
        self.image_tags_label = QLabel("")
        self.image_tags_label.setAlignment(Qt.AlignCenter)
        self.image_tags_label.setStyleSheet(f"""
            color: #2196F3;
            padding: 5px;
            font-size: {self.font_sizes['small']}px;
        """)
        self.image_tags_label.setWordWrap(True)
        self.image_tags_label.hide()
        layout.addWidget(self.image_tags_label)

        return page
    
    def create_text_page(self) -> QWidget:
        """创建文本预览页面"""
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setContentsMargins(0, 0, 0, 0)

        self.text_edit = QTextEdit()
        self.text_edit.setReadOnly(True)
        self.text_edit.setFont(QFont("Consolas", self.font_sizes['normal']))
        self.text_edit.setStyleSheet("""
            QTextEdit {
                border: none;
                background-color: #fafafa;
                padding: 10px;
            }
        """)

        layout.addWidget(self.text_edit)

        return page

    def create_audio_page(self) -> QWidget:
        """创建音频预览页面"""
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setAlignment(Qt.AlignCenter)

        icon_label = QLabel("🎵")
        icon_label.setStyleSheet(f"font-size: {self.font_sizes['icon_large']}px;")
        icon_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(icon_label)

        label = QLabel("音频文件")
        label.setStyleSheet(f"font-size: {self.font_sizes['icon_medium']}px; color: #333;")
        label.setAlignment(Qt.AlignCenter)
        layout.addWidget(label)

        desc_label = QLabel("音频预览需要额外的播放器\n请使用外部程序播放")
        desc_label.setStyleSheet("color: #666;")
        desc_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(desc_label)

        open_btn = QPushButton("使用默认程序播放")
        open_btn.setMinimumHeight(self.window_sizes['button_height'])
        open_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: #3498db;
                color: white;
                padding: 10px 20px;
                border-radius: 4px;
                font-size: {self.font_sizes['button']}px;
            }}
            QPushButton:hover {{
                background-color: #2980b9;
            }}
        """)
        open_btn.clicked.connect(self.open_current_file)
        layout.addWidget(open_btn, alignment=Qt.AlignCenter)

        return page

    def create_unsupported_page(self) -> QWidget:
        """创建不支持页面"""
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setAlignment(Qt.AlignCenter)

        icon_label = QLabel("📎")
        icon_label.setStyleSheet(f"font-size: {self.font_sizes['icon_large']}px;")
        icon_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(icon_label)

        label = QLabel("不支持的文件格式")
        label.setStyleSheet(f"font-size: {self.font_sizes['icon_medium']}px; color: #333;")
        label.setAlignment(Qt.AlignCenter)
        layout.addWidget(label)

        self.unsupported_desc_label = QLabel("该文件类型暂不支持预览\n请使用外部程序打开")
        self.unsupported_desc_label.setStyleSheet("color: #666;")
        self.unsupported_desc_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(self.unsupported_desc_label)

        open_btn = QPushButton("使用默认程序打开")
        open_btn.setMinimumHeight(self.window_sizes['button_height'])
        open_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: #95a5a6;
                color: white;
                padding: 10px 20px;
                border-radius: 4px;
                font-size: {self.font_sizes['button']}px;
            }}
            QPushButton:hover {{
                background-color: #7f8c8d;
            }}
        """)
        open_btn.clicked.connect(self.open_current_file)
        layout.addWidget(open_btn, alignment=Qt.AlignCenter)

        return page

    def create_loading_page(self) -> QWidget:
        """创建加载中页面"""
        page = QWidget()
        layout = QVBoxLayout(page)
        layout.setAlignment(Qt.AlignCenter)

        self.loading_label = QLabel("正在加载预览...")
        self.loading_label.setStyleSheet(f"""
            color: #666;
            font-size: {self.font_sizes['title']}px;
        """)
        self.loading_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(self.loading_label)

        self.loading_bar = QProgressBar()
        self.loading_bar.setRange(0, 0)  # 无限进度
        self.loading_bar.setMaximumWidth(200)
        layout.addWidget(self.loading_bar, alignment=Qt.AlignCenter)

        return page
    
    def preview_file(self, file_path: str, file_id: Optional[int] = None,
                     metadata: dict = None, category_system_name: Optional[str] = None):
        """预览文件

        Args:
            file_path: 文件路径
            file_id: 文件ID
            metadata: 文件元数据（包含caption、tags等）
            category_system_name: 当前类别体系名称
        """
        if not file_path or not os.path.exists(file_path):
            self.stack.setCurrentIndex(0)  # 空白页面
            return

        self.current_file = file_path
        self.current_metadata = metadata
        self.current_file_id = file_id
        self.current_category_system = category_system_name

        # 更新文件信息
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        size_str = self.format_size(file_size)
        self.file_info_label.setText(f"{file_name} ({size_str})")

        # 显示加载页面
        self.stack.setCurrentIndex(9)  # 加载中页面

        # 启动后台加载
        if self.preview_worker and self.preview_worker.isRunning():
            self.preview_worker.wait()

        self.preview_worker = PreviewWorker(file_path)
        self.preview_worker.preview_ready.connect(self.on_preview_ready)
        self.preview_worker.preview_error.connect(self.on_preview_error)
        self.preview_worker.start()
    
    def on_preview_ready(self, data, preview_type: str):
        """预览加载完成"""
        if preview_type == 'image':
            self.show_image_preview(data)
        elif preview_type == 'text':
            self.show_text_preview(data)
        elif preview_type == 'pdf':
            self.show_pdf_preview()
        elif preview_type == 'word':
            self.show_word_preview()
        elif preview_type == 'excel':
            self.show_excel_preview()
        elif preview_type == 'ppt':
            self.show_ppt_preview()
        elif preview_type == 'audio':
            self.stack.setCurrentIndex(7)  # 音频页面
        elif preview_type == 'shortcut':
            self.show_shortcut_preview()
        else:
            self.stack.setCurrentIndex(8)  # 不支持页面

        # 显示语义块信息
        self._show_semantic_blocks_info()
    
    def on_preview_error(self, error_msg: str):
        """预览加载错误"""
        self.loading_label.setText(f"加载失败: {error_msg}")
        self.loading_bar.hide()

        # 显示语义块信息（即使预览失败也显示）
        self._show_semantic_blocks_info()

    def _show_semantic_blocks_info(self):
        """显示语义块信息"""
        if not self.current_file_id or not self.current_category_system:
            self.semantic_blocks_label.setVisible(False)
            return

        from database import get_db_manager
        db_manager = get_db_manager()

        semantic_blocks = db_manager.get_semantic_blocks_by_file(self.current_file_id)

        if not semantic_blocks:
            self.semantic_blocks_label.setVisible(False)
            return

        # 构建显示文本
        info_lines = ["语义块信息:"]
        for block in semantic_blocks:
            block_id = block.semantic_block_id
            metadata = block.metadata

            if metadata and self.current_category_system in metadata:
                similarities = metadata[self.current_category_system]
                # 格式化相似度显示，显示前3个最相似的类别
                top_cats = sorted(similarities.items(), key=lambda x: -x[1])[:3]
                sim_str = ", ".join([f"{cat}:{sim:.1%}" for cat, sim in top_cats])
                info_lines.append(f"  {block_id[:20]}... | {sim_str}")
            else:
                info_lines.append(f"  {block_id[:20]}... | 无分类数据")

        # 显示在信息区域
        self.semantic_blocks_label.setText("\n".join(info_lines))
        self.semantic_blocks_label.setVisible(True)

    def show_image_preview(self, pixmap: QPixmap):
        """显示图片预览"""
        # 缩放图片以适应窗口
        max_size = 800
        if pixmap.width() > max_size or pixmap.height() > max_size:
            pixmap = pixmap.scaled(max_size, max_size,
                                  Qt.KeepAspectRatio,
                                  Qt.SmoothTransformation)

        self.image_label.setPixmap(pixmap)

        # 构建信息文本
        info_parts = [f"{pixmap.width()} x {pixmap.height()} px"]

        # 添加元数据信息
        if hasattr(self, 'current_metadata') and self.current_metadata:
            metadata = self.current_metadata
            # 拍摄时间
            if metadata.get('capture_time_extracted'):
                info_parts.append(f"拍摄: {metadata['capture_time_extracted']}")
            elif metadata.get('capture_time_from_filename'):
                info_parts.append(f"时间: {metadata['capture_time_from_filename']}")
            # 地点
            if metadata.get('location_info'):
                info_parts.append(f"地点: {metadata['location_info']}")

        self.image_info_label.setText(" | ".join(info_parts))

        # 显示Caption
        if hasattr(self, 'current_metadata') and self.current_metadata and self.current_metadata.get('caption'):
            self.image_caption_label.setText(f"📷 {self.current_metadata['caption']}")
            self.image_caption_label.show()
        else:
            self.image_caption_label.hide()

        # 显示Tags
        if hasattr(self, 'current_metadata') and self.current_metadata and self.current_metadata.get('tags'):
            tags = self.current_metadata['tags']
            tags_text = " ".join([f"#{tag}" for tag in tags if tag])
            self.image_tags_label.setText(f"🏷️ {tags_text}")
            self.image_tags_label.show()
        else:
            self.image_tags_label.hide()

        self.stack.setCurrentIndex(1)  # 图片页面
    
    def show_text_preview(self, content: str):
        """显示文本预览"""
        self.text_edit.setPlainText(content)
        self.stack.setCurrentIndex(2)  # 文本页面
    
    def show_pdf_preview(self):
        """显示PDF预览"""
        if PDF_SUPPORT:
            try:
                success = self.pdf_widget.load_pdf(self.current_file)
                if success:
                    self.stack.setCurrentIndex(3)  # PDF页面
                else:
                    self.stack.setCurrentIndex(8)  # 不支持页面
            except Exception:
                self.stack.setCurrentIndex(8)  # 不支持页面
        else:
            self.stack.setCurrentIndex(8)  # 不支持页面
    
    def show_word_preview(self):
        """显示Word预览"""
        file_lower = self.current_file.lower()
        if DOCX_SUPPORT and file_lower.endswith('.docx'):
            try:
                success = self.word_widget.load_document(self.current_file)
                if success:
                    self.stack.setCurrentIndex(4)  # Word页面
                else:
                    self.stack.setCurrentIndex(8)  # 不支持页面
            except Exception:
                self.stack.setCurrentIndex(8)  # 不支持页面
        elif file_lower.endswith('.doc'):
            # 旧版Word格式不支持，显示特定提示
            self.unsupported_desc_label.setText("旧版Word格式(.doc)暂不支持预览\n请转换为.docx格式或使用外部程序打开")
            self.stack.setCurrentIndex(8)  # 不支持页面
        else:
            self.unsupported_desc_label.setText("该文件类型暂不支持预览\n请使用外部程序打开")
            self.stack.setCurrentIndex(8)  # 不支持页面
    
    def show_excel_preview(self):
        """显示Excel预览"""
        file_lower = self.current_file.lower()
        if XLSX_SUPPORT and file_lower.endswith('.xlsx'):
            try:
                success = self.excel_widget.load_workbook(self.current_file)
                if success:
                    self.stack.setCurrentIndex(5)  # Excel页面
                else:
                    self.stack.setCurrentIndex(8)  # 不支持页面
            except Exception:
                self.stack.setCurrentIndex(8)  # 不支持页面
        elif file_lower.endswith('.xls'):
            # 旧版Excel格式不支持，显示特定提示
            self.unsupported_desc_label.setText("旧版Excel格式(.xls)暂不支持预览\n请转换为.xlsx格式或使用外部程序打开")
            self.stack.setCurrentIndex(8)  # 不支持页面
        else:
            self.unsupported_desc_label.setText("该文件类型暂不支持预览\n请使用外部程序打开")
            self.stack.setCurrentIndex(8)  # 不支持页面
    
    def show_ppt_preview(self):
        """显示PowerPoint预览"""
        file_lower = self.current_file.lower()
        if PPTX_SUPPORT and file_lower.endswith('.pptx'):
            try:
                success = self.ppt_widget.load_presentation(self.current_file)
                if success:
                    self.stack.setCurrentIndex(6)  # PPT页面
                else:
                    self.stack.setCurrentIndex(8)  # 不支持页面
            except Exception:
                self.stack.setCurrentIndex(8)  # 不支持页面
        elif file_lower.endswith('.ppt'):
            # 旧版PowerPoint格式不支持，显示特定提示
            self.unsupported_desc_label.setText("旧版PowerPoint格式(.ppt)暂不支持预览\n请转换为.pptx格式或使用外部程序打开")
            self.stack.setCurrentIndex(8)  # 不支持页面
        else:
            self.unsupported_desc_label.setText("该文件类型暂不支持预览\n请使用外部程序打开")
            self.stack.setCurrentIndex(8)  # 不支持页面
    
    def show_shortcut_preview(self):
        """显示快捷方式预览"""
        # 显示特定提示
        self.unsupported_desc_label.setText("快捷方式(.lnk)暂不支持预览\n请使用外部程序打开")
        self.stack.setCurrentIndex(8)  # 不支持页面
    
    def open_current_file(self):
        """使用系统默认程序打开当前文件"""
        if self.current_file and os.path.exists(self.current_file):
            import subprocess
            import platform
            
            try:
                if platform.system() == 'Windows':
                    os.startfile(self.current_file)
                elif platform.system() == 'Darwin':  # macOS
                    subprocess.run(['open', self.current_file])
                else:  # Linux
                    subprocess.run(['xdg-open', self.current_file])
            except Exception as e:
                QMessageBox.warning(self, "打开失败", f"无法打开文件: {e}")
    
    def format_size(self, size: int) -> str:
        """格式化文件大小"""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size < 1024:
                return f"{size:.1f} {unit}"
            size /= 1024
        return f"{size:.1f} TB"
    
    def clear_preview(self):
        """清除预览"""
        self.current_file = None
        self.file_info_label.clear()
        self.stack.setCurrentIndex(0)  # 空白页面
        
        # 关闭PDF文档
        self.pdf_widget.close_document()